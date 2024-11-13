import os
import tkinter as tk
from tkinter import filedialog
from datetime import datetime
import pandas as pd
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
import openpyxl
from openpyxl.styles import Alignment

# Ruta del escritorio y archivo de log
desktop_path = os.path.join(os.path.expanduser("~"), "Desktop")
log_file = os.path.join(desktop_path, "renombrados_log.xlsx")

# Listas para búsqueda de palabras clave
metodologias = {
    "Back Data": ["back", "data", "histórico", "antecedentes","Back Data"],
    "Canasto Alcohólico": ["canasto", "alcohólico", "bebidas", "alcohol","Canasto Alcohólico"],
    "Coberturas": ["cobertura", "coverage","Coberturas"],
    "Golden Stores": ["golden", "stores", "Golden Stores"],
    "Consumer Pulse Survey": ["consumer", "survey", "encuesta","Consumer Pulse Survey","pulse"],
    "Universos de Tiendas": ["universo", "tiendas", "Universos de Tiendas","Universos"],
    "Indicadores de Desempeño": ["indicadores", "desempeño", "business performance","Indicadores de Desempeño","Desempeño"],
    "Tracking Distribution": ["tracking", "distribution","Tracking Distribution","Auditoría"],
    "Estudio de precios": ["estudio", "precios", "ventas","Estudio de precios"],
    "Brand Board": ["brand", "brand board", "brandboards","Brand Board"],
    "Category Overview": ["Overview", "overview sport drinks", "overview frijoles", "Category Overview"],
    "Reporte BDP": ["reporte", "bdp","Reporte BDP"],
    "Brand power": ["kantar", "mercaplan", "kantar mercaplan","Brand power"],
    "Brand Book": ["brand", "book","Brand Book"],
    "Análisis": ["Análisis","análisis"],
    "Diagnóstico de Categorías": ["Diagnóstico","Diagnóstico de Categorías"],
    "Convivencias entre Categorías": ["Convivencias entre Categorías"],
    "Webinar": ["Webinar","webinar"],
    "Post Lanzamiento": [ "Post Lanzamiento","Launch","post launch"],
    "Estudio de Usos y Hábitos": ["Hábitos", "Estudio de Usos y Hábitos"],
    "Evaluación de Empaque": ["Empaque","Evaluación de Empaque"],
    "Post Test Publicitario": ["Post Test Publicitario","Test","Prueba","Publicitario"],
    "Pre Test Publicitario": ["Pre Test Publicitario","Test","Prueba","Publicitario"],
    "Evaluación de Concepto": ["Concepto","Evaluación de Concepto"],
    "Esencia de Marca": ["Esencia de Marca","Esencia","Marca"],
    "Tracker de Canal": ["Tracker", "Canal","Tracker de Canal"],
    "Lealtad de Clientes": ["Lealtad","Clientes","Lealtad de Clientes"],
    "Evaluación de Producto": ["Producto","Evaluación de Producto"],
    "Shopper Insights": ["Shopper","Shopper Insights","Insights"]

}

categorias = {
    "Cervezas": ["cerveza", "beer","Cervezas","Beer","Cerveza","Guaro"],
    "RTD´s": ["rtd", "ready-to-drink", "bas", "mezclados","BAS", "RTD´s"],
    "Té Frío": ["té", "cold tea", "tropical","Té Frío","Té","Frío","TES"],
    "Energéticas": ["energético", "energy drinks","Energéticas","Maxxx Energy"],
    "Gaseosas": ["gaseosa", "soda", "csd's","gaseosas","Gaseosas","Pepsi","Milory"],
    "Isotónicas": ["isotónica", "isotonic", "sport drinks","Isotónicas"],
    "Frijoles": ["frijol", "beans","Frijoles","Beans","Ducal","frijoles"],
    "Ketchup": ["ketchup", "catsup", "kern's", "Ketchup","Salsas"],
    "Néctares": ["néctares", "nectar","Néctares"],
    "Vinos y Destilados": ["vino", "wines","Vinos y Destilados","Destilados","Vinos"],
    "Categorías Emergentes": ["emergente","Categorías Emergentes"],
    "Multicategoría": ["multi-category","Multicategoría"],
    "Nutrition": ["nutrition","Nutrición","Nutrition"],
    "Bebidas Alcohólicas": ["alcohólica","Bebidas Alcohólicas"],
    "Snacks": ["snack","Snacks"],
    "RFBs": ["rfb","RFBs"],
    "Sueros": ["suero", "gatorlyte","Sueros"],
    "Café": ["coffee","Café"],
    "Bebidas funcionales": ["funcional","Bebidas funcionales"],
    "Arroz RTE": ["arroz", "rte","Arroz RTE"]
}

regiones = [
    "Multipaís", "Estados Unidos", "México", "Guatemala", "El Salvador",
    "Honduras", "Nicaragua", "Costa Rica", "Panamá",
    "República Dominicana", "LATAM", "Europa", "Centroamérica"
]

eje_tematico = {
    "Entrenamientos": ["training", "entrenamiento", "capacitacion", "Entrenamientos","Entrenamiento"],
    "Publicidad y Comunicación": ["publicidad", "comunicación", "campaña", "Publicidad y Comunicación"],
    "Desempeño de negocio": ["desempeño", "performance", "ventas","Desempeño de negocio","Indicadores","Desempeño"],
    "Innovaciones": ["innovación", "tecnología","Innovaciones","Nuevo","Nueva"],
    "Shoppers y Canales": ["shoppers", "canales","Shoppers y Canales","Supermercados","Compradores","Tienda","Tiendas"],
    "Portafolio": ["portafolio", "productos","Portafolio"],
    "Diagnósticos": ["diagnóstico", "evaluacion","Diagnósticos","Diagnóstico"],
    "Macroeconómicos": ["macroeconomía", "finanzas","Macroeconómicos", "Contexto Macroeconómico"],
    "Macrotendencias": ["macrotendencias", "mercado","Macrotendencias"]
}

marcas = {
    "Imperial": ["imperial", "cerveza imperial","Imperial","IMPERIAL","Imperial Silver"],
    "Pilsen": ["pilsen", "cerveza pilsen","Pilsen"],
    "Bavaria": ["bavaria", "cerveza bavaria","Bavaria"],
    "Heineken": ["heineken", "cerveza heineken","Heineken"],
    "Rock Limon": ["rock limon", "limon soda","Rock Limon","Rock"],
    "Bohemia": ["bohemia", "cerveza bohemia","Bohemia"],
    "Adán & Eva": ["adán", "eva","Adán & Eva"],
    "Bamboo": ["bamboo", "bambu","Bamboo"],
    "Smirnoff Ice": ["smirnoff ice", "smirnoff","Smirnoff Ice"],
    "Maxxx Energy": ["maxxx energy", "maxxx","energy","Maxxx Energy"],
    "Jet": ["jet", "Jet"],
    "Tropical": ["tropical", "té frío","Tropical"],
    "Pepsi": ["pepsi", "pepsico","Gaseosas","Pepsi"],
    "Kern´s": ["kern's", "ketchup","Kern´s"],
    "Ducal": ["ducal", "frijoles","Ducal"],
    "Gatorlyte": ["gatorlyte", "sueros", "Gatorlyte"],
    "Gatorade": ["gatorade", "Gatorade"],
    "Multimarca": ["multimarca","Multimarca"],
    "Competencia": ["competencia","Competencia"],
    "Sol": ["sol","Sol"],
    "Vida": ["vida","Vida"]
}

def buscar_mejor_ajuste(nombre, lista_palabras):
    nombre = nombre.lower()
    for item, palabras in lista_palabras.items():
        if any(palabra in nombre for palabra in palabras):
            return item
    return "Desconocido"

def extraer_fecha_de_archivo(ruta_archivo):
    # Aquí se añade lógica para PDF, DOCX, PPTX
    try:
        ext = os.path.splitext(ruta_archivo)[1].lower()
        if ext == ".pdf":
            reader = PdfReader(ruta_archivo)
            for page in reader.pages:
                texto = page.extract_text().lower()
                for mes, nombre_mes in [
                    ("ene", "ene"), ("feb", "feb"), ("mar", "mar"), 
                    ("abr", "abr"), ("may", "may"), ("jun", "jun"), 
                    ("jul", "jul"), ("ago", "ago"), ("sep", "sep"), 
                    ("oct", "oct"), ("nov", "nov"), ("dic", "dic")
                ]:
                    if mes in texto:
                        año = ''.join(filter(str.isdigit, texto))
                        if len(año) == 4:  # Año en formato de 4 dígitos
                            return f"{nombre_mes},{año}"
        elif ext == ".docx":
            doc = Document(ruta_archivo)
            for paragraph in doc.paragraphs:
                texto = paragraph.text.lower()
                for mes, nombre_mes in [
                    ("ene", "ene"), ("feb", "feb"), ("mar", "mar"), 
                    ("abr", "abr"), ("may", "may"), ("jun", "jun"), 
                    ("jul", "jul"), ("ago", "ago"), ("sep", "sep"), 
                    ("oct", "oct"), ("nov", "nov"), ("dic", "dic")
                ]:
                    if mes in texto:
                        año = ''.join(filter(str.isdigit, texto))
                        if len(año) == 4:
                            return f"{nombre_mes},{año}"
        elif ext == ".pptx":
            ppt = Presentation(ruta_archivo)
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        texto = shape.text.lower()
                        for mes, nombre_mes in [
                            ("ene", "ene"), ("feb", "feb"), ("mar", "mar"), 
                            ("abr", "abr"), ("may", "may"), ("jun", "jun"), 
                            ("jul", "jul"), ("ago", "ago"), ("sep", "sep"), 
                            ("oct", "oct"), ("nov", "nov"), ("dic", "dic")
                        ]:
                            if mes in texto:
                                año = ''.join(filter(str.isdigit, texto))
                                if len(año) == 4:
                                    return f"{nombre_mes},{año}"
    except Exception as e:
        print(f"Error al extraer fecha del archivo {ruta_archivo}: {e}")
    return "Fecha desconocida"

def obtener_region(ruta_archivo):
    for region in regiones:
        if region.lower() in ruta_archivo.lower():
            return region
    return "Región desconocida"

def renombrar_archivo(ruta_archivo):
    nombre_original = os.path.splitext(os.path.basename(ruta_archivo))[0]
    metodologia = buscar_mejor_ajuste(nombre_original, metodologias)
    categoria = buscar_mejor_ajuste(nombre_original, categorias)
    fecha = extraer_fecha_de_archivo(ruta_archivo)
    region = obtener_region(ruta_archivo)
    eje = buscar_mejor_ajuste(nombre_original, eje_tematico)
    marca = buscar_mejor_ajuste(nombre_original, marcas)
    nuevo_nombre = f"{metodologia}|{categoria}|{region}|{fecha}"
    nuevo_nombre_log = [nombre_original, metodologia, categoria, region, fecha, eje, marca, nuevo_nombre]

    columnas = ["Nombre Original", "Metodología", "Categoría", "Región", "Fecha", "Eje Temático", "Marca", "Renombrado"]

    if os.path.exists(log_file):
        wb = openpyxl.load_workbook(log_file)
        ws = wb.active
        ws.append(nuevo_nombre_log)
        wb.save(log_file)
    else:
        df_log = pd.DataFrame([nuevo_nombre_log], columns=columnas)
        df_log.to_excel(log_file, index=False)

def seleccionar_archivos():
    archivos = filedialog.askopenfilenames(title="Seleccionar archivos", filetypes=(("All Files", "*.*"),))
    for archivo in archivos:
        renombrar_archivo(archivo)

root = tk.Tk()
root.title("Renombrar Archivos")
root.geometry("400x200")

seleccionar_btn = tk.Button(root, text="Seleccionar Archivos", command=seleccionar_archivos)
seleccionar_btn.pack(pady=20)

root.mainloop()
